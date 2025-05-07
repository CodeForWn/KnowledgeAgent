import datetime
import logging
from pymongo import MongoClient, errors
import sys
import os
import json
import pdfplumber
from docx import Document
from pptx import Presentation
import random
import re
import time

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
from config.KMC_config import Config

class KMCMongoDBHandler:
    def __init__(self, config):
        try:
            # 从配置中提取MongoDB相关配置
            self.mongodb_host = config.mongodb_host
            self.mongodb_port = config.mongodb_port
            self.mongodb_database = config.mongodb_database
            self.mongodb_username = getattr(config, "mongodb_username", None)
            self.mongodb_password = getattr(config, "mongodb_password", None)
            self.config = config
            self.logger = self.config.logger

            # 建立MongoDB连接
            if self.mongodb_username and self.mongodb_password:
                self.client = MongoClient(
                    host=self.mongodb_host,
                    port=self.mongodb_port,
                    username=self.mongodb_username,
                    password=self.mongodb_password
                )
            else:
                self.client = MongoClient(host=self.mongodb_host, port=self.mongodb_port)
            self.db = self.client[self.mongodb_database]
            self.logger.info("成功连接到 MongoDB {}:{}".format(self.mongodb_host, self.mongodb_port))
        except errors.ConnectionFailure as e:
            self.logger.error("MongoDB连接失败: {}".format(e))
        except Exception as e:
            self.logger.error("连接MongoDB时出现错误: {}".format(e))

    def insert_document(self, collection_name, document):
        """向指定集合中插入单个文档"""
        try:
            collection = self.db[collection_name]
            result = collection.insert_one(document)
            self.logger.info("向集合 {} 插入文档成功，ID: {}".format(collection_name, result.inserted_id))
            return result.inserted_id
        except Exception as e:
            self.logger.error("向集合 {} 插入文档失败: {}".format(collection_name, e))
            return None

    def generate_docid(self):
        """
        生成一个唯一的 docID，格式为三个4位数字拼接，如 "1234-5678-9012"
        """
        return "{}-{}-{}".format(
            str(random.randint(0, 9999)).zfill(4),
            str(random.randint(0, 9999)).zfill(4),
            str(random.randint(0, 9999)).zfill(4)
        )

    def insert_question(self, collection_name, question_data):
        """向题库中插入一个新题目"""
        try:
            result = self.db[collection_name].insert_one(question_data)
            self.logger.info(f"插入试题 docID={question_data['docID']}，插入成功")
            return result.inserted_id
        except Exception as e:
            self.logger.error(f"插入试题失败 docID={question_data['docID']}：{e}")
            return None

    def close(self):
        """关闭MongoDB连接"""
        try:
            self.client.close()
            self.logger.info("关闭MongoDB连接成功")
        except Exception as e:
            self.logger.error("关闭MongoDB连接时出错: {}".format(e))

    def get_resource_by_docID(self, docID, collection_name="geo_documents"):
        """
        根据 docID 查询 MongoDB 中的资源详细信息
        """
        try:
            collection = self.db[collection_name]
            resource_info = collection.find_one({"docID": docID})
            if resource_info:
                # 将 MongoDB 的 _id 转换为字符串，避免序列化问题
                if "_id" in resource_info:
                    resource_info["_id"] = str(resource_info["_id"])
            return resource_info
        except Exception as e:
            self.logger.error("查询 docID {} 出错: {}".format(docID, e))
            return None

    def get_folder_id_by_docID(self, docID):
        """
        根据 docID 查询 MongoDB 中的资源详细信息。
        先在 geo_documents 查找，如果找不到，再去 edu_question 查找。
        """
        try:
            # 先查 geo_documents
            collection = self.db["geo_documents"]
            resource_info = collection.find_one({"docID": docID})

            if resource_info:
                if "_id" in resource_info:
                    resource_info["_id"] = str(resource_info["_id"])
                return resource_info

            # 如果在 geo_documents 没找到，查 edu_question
            collection = self.db["edu_question"]
            resource_info = collection.find_one({"docID": docID})

            if resource_info:
                if "_id" in resource_info:
                    resource_info["_id"] = str(resource_info["_id"])
                return resource_info

            # 都找不到
            return None

        except Exception as e:
            self.logger.error(f"查询 docID {docID} 出错: {e}")
            return None

    @staticmethod
    def read_pdf(file_path):
        """
        读取 PDF 文件的全文内容
        """
        try:
            content = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n"
            return content
        except Exception as e:
            return f"读取PDF文件出错: {e}"

    @staticmethod
    def read_docx(file_path):
        """
        读取 Word 文件的全文内容
        """
        try:
            doc = Document(file_path)
            content = "\n".join([p.text for p in doc.paragraphs])
            return content
        except Exception as e:
            return f"读取Word文件出错: {e}"

    @staticmethod
    def read_txt(file_path):
        """
        读取 TXT 文件的全文内容
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            return f"读取TXT文件出错: {e}"

    @staticmethod
    def read_ppt(file_path):
        """
        读取 PPT 或 PPTX 文件的全文内容
        """
        try:
            prs = Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"
            return content
        except Exception as e:
            return f"读取PPT文件出错: {e}"

    @staticmethod
    def read_file_content(file_path):
        """
        根据文件后缀选择合适的方式读取全文内容
        """
        if file_path.endswith(".pdf"):
            return KMCMongoDBHandler.read_pdf(file_path)
        elif file_path.endswith(".docx"):
            return KMCMongoDBHandler.read_docx(file_path)
        elif file_path.endswith(".txt"):
            return KMCMongoDBHandler.read_txt(file_path)
        elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
            return KMCMongoDBHandler.read_ppt(file_path)
        else:
            return "不支持的文件格式。"

    def get_resources_full_text(self, resources, collection_name="geo_documents"):
        """
        针对从Neo4j接口返回的资源列表（包含docID、file_name、resource_type），
        查询MongoDB中对应的资源详细信息（包括文件路径），然后读取该文件的全文内容，
        并将全文内容合并到资源字典中返回。
        """
        detailed_resources = []
        for res in resources:
            docID = res.get("docID")
            resource_info = self.get_resource_by_docID(docID, collection_name)
            if resource_info:
                file_path = resource_info.get("file_path", "")
                full_text = self.read_file_content(file_path)
                res["full_text"] = full_text
                # 合并 MongoDB 中其它详细信息，并确保 ObjectId 转为字符串
                for k, v in resource_info.items():
                    if k == "_id":
                        res[k] = str(v)
                    else:
                        res[k] = v
            else:
                res["full_text"] = "未在MongoDB中找到该资源信息。"
            detailed_resources.append(res)
        return detailed_resources


    def filter_documents(self, filters: dict, collection_name="geo_documents"):
        """根据多条件筛选资源库数据（resource_type != 试题）"""
        query = {
            "resource_type": {"$ne": "试题"}
        }
        if filters.get("kb_id"):
            query["kb_id"] = filters["kb_id"]
        if filters.get("folder_id"):
            query["folder_id"] = filters["folder_id"]
        if filters.get("file_name"):
            # ✅ 改成模糊匹配（多个关键词之间是 OR 关系）
            file_name_patterns = filters["file_name"]
            if isinstance(file_name_patterns, list) and file_name_patterns:
                query["$or"] = [{"file_name": {"$regex": fn, "$options": "i"}} for fn in file_name_patterns]
        if filters.get("docID"):
            query["docID"] = filters["docID"]  # ✅ 新增
        if filters.get("status"):
            query["status"] = {"$in": filters["status"]}
        if filters.get("knowledge_point"):
            kp_values = filters["knowledge_point"]
            if isinstance(kp_values, list) and kp_values:
                query["knowledge_point"] = {"$in": kp_values}

        return list(self.db[collection_name].find(query).sort("created_at", -1))

    def filter_questions(self, filters: dict, collection_name="edu_question"):
        """根据多条件筛选题库数据（resource_type = 试题）"""
        query = {
            "resource_type": "试题"
        }
        if filters.get("kb_id"):
            query["kb_id"] = filters["kb_id"]
        if filters.get("folder_id"):
            query["folder_id"] = filters["folder_id"]
        if filters.get("type"):
            query["type"] = {"$in": filters["type"]}
        if filters.get("diff_level"):
            query["diff_level"] = {"$in": filters["diff_level"]}
        if filters.get("docID"):
            query["docID"] = filters["docID"]  # ✅ 新增
        if filters.get("status"):
            query["status"] = {"$in": filters["status"]}
        if filters.get("knowledge_point"):
            query["knowledge_point"] = {"$in": filters["knowledge_point"]}

        return list(self.db[collection_name].find(query).sort("created_at", -1))

    def update_document_by_id(self, doc_id, update_fields, collection_name="geo_documents"):
        try:
            result = self.db[collection_name].update_one({"docID": doc_id}, {"$set": update_fields})
            self.logger.info(f"更新文档 docID={doc_id}，修改字段：{update_fields}")
            return result
        except Exception as e:
            self.logger.error(f"更新文档失败 docID={doc_id}：{e}")
            return None

    def update_question_by_id(self, doc_id, update_fields, collection_name="edu_question"):
        try:
            result = self.db[collection_name].update_one({"docID": doc_id}, {"$set": update_fields})
            self.logger.info(f"更新试题 docID={doc_id}，修改字段：{update_fields}")
            return result
        except Exception as e:
            self.logger.error(f"更新试题失败 docID={doc_id}：{e}")
            return None

    def delete_document_by_id(self, doc_id, collection_name="geo_documents"):
        try:
            result = self.db[collection_name].delete_one({"docID": doc_id})
            self.logger.info(f"删除资源 docID={doc_id}，删除数：{result.deleted_count}")
            return result
        except Exception as e:
            self.logger.error(f"删除资源失败 docID={doc_id}：{e}")
            return None

    def delete_question_by_id(self, doc_id, collection_name="edu_question"):
        try:
            result = self.db[collection_name].delete_one({"docID": doc_id})
            self.logger.info(f"删除试题 docID={doc_id}，删除数：{result.deleted_count}")
            return result
        except Exception as e:
            self.logger.error(f"删除试题失败 docID={doc_id}：{e}")
            return None


    # ✅ 新增：列出 geo_documents 中 folder_id 为空的 docID
    def get_documents_with_empty_folder_id(self, collection_name="geo_documents"):
        try:
            collection = self.db[collection_name]
            query = {
                "$or": [
                    {"folder_id": {"$exists": False}},
                    {"folder_id": ""},
                    {"folder_id": None}
                ]
            }
            documents = collection.find(query, {"docID": 1, "_id": 0})
            empty_folder_docIDs = [doc["docID"] for doc in documents if "docID" in doc]
            self.logger.info(f"找到 {len(empty_folder_docIDs)} 个 folder_id 为空的文档")
            return empty_folder_docIDs
        except Exception as e:
            self.logger.error(f"查询 folder_id 为空的文档出错: {e}")
            return []

# if __name__ == '__main__':
#     # 示例用法
#     config = Config()  # 假设 Config 类中包含mongodb_host、mongodb_port、mongodb_database、logger等配置项
#     config.load_config()
#     mongo_handler = KMCMongoDBHandler(config)
#     # ✅ 调用新增的方法
#     empty_folder_docIDs = mongo_handler.get_documents_with_empty_folder_id()
#     print(f"共有 {len(empty_folder_docIDs)} 个空 folder_id 的文档")
#     for doc_id in empty_folder_docIDs:
#         print(doc_id)
#
#     # 示例：根据 docID 查询资源信息并读取全文内容
#     test_docID = "0033-8784-3716"
#     resource = mongo_handler.get_resource_by_docID(test_docID)
#     print(resource)
#     if resource:
#         file_path = resource.get("file_path", "")
#         full_text = mongo_handler.read_file_content(file_path)
#         resource["full_text"] = full_text
#
#         # print(resource)
#     else:
#         print("未找到对应的资源信息。")
#
#     # 假设有一组从Neo4j返回的资源列表
#     resources = [
#         {
#             "docID": "0033-8784-3716",
#             "file_name": "第一单元地球自转部分.pdf",
#             "resource_type": "教材"
#         },
#         {
#             "docID": "5737-1173-5069",
#             "file_name": "主题1 地球自转 第一课时.pptx",
#             "resource_type": "课件"
#         }
#     ]
#     detailed_resources = mongo_handler.get_resources_full_text(resources)
#     # print(json.dumps(detailed_resources, indent=2, ensure_ascii=False, default=str))
#
#     mongo_handler.close()